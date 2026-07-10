"""`docs/uir.schema.json` is the published contract. It must match the model.

The schema is generated from `UIRV1` by `scripts/export_uir_json_schema.py`,
but nothing regenerated it when PLAN §17 widened the models. It drifted three
ways, and *every* UIR the pipeline emitted -- including a plain PDF -- failed
validation against the file we hand downstream consumers:

    source.format   was `const: "PDF"`; the model now allows DOCX/PPTX/...
    source.route    rejected as an additional property
    metadata.format rejected as an additional property

Nothing caught it because Pydantic validates documents against the *model*,
never against the exported artifact.
"""
from __future__ import annotations

import json
from pathlib import Path

import pytest

jsonschema = pytest.importorskip("jsonschema")

from uir_pipeline.uir_schema import schema_json_dict  # noqa: E402

_SCHEMA_PATH = Path(__file__).resolve().parent.parent / "docs" / "uir.schema.json"


def _on_disk() -> dict:
    return json.loads(_SCHEMA_PATH.read_text(encoding="utf-8"))


def test_exported_schema_exists():
    assert _SCHEMA_PATH.is_file()


def test_exported_schema_matches_the_model():
    """Regenerate with `python scripts/export_uir_json_schema.py` when this fails."""
    assert _on_disk() == schema_json_dict(), (
        "docs/uir.schema.json is stale; run scripts/export_uir_json_schema.py"
    )


def test_schema_allows_the_multi_format_source_fields():
    """PLAN §17 widened Source.format and added Source.route / Metadata.format."""
    schema = _on_disk()
    source = schema["$defs"]["Source"]["properties"]
    assert "route" in source, "Source.route missing from the published schema"
    assert "format" in source
    # No longer pinned to PDF.
    assert source["format"].get("const") != "PDF"
    assert "format" in schema["$defs"]["Metadata"]["properties"]


def _uir_from(tmp_path: Path, source_path: Path):
    from uir_pipeline.pipeline import run

    out = tmp_path / "out"
    result = run(source_path, output_dir=out, skip_weaviate=True, with_embeddings=False)
    return json.loads(Path(result.out_path).read_text(encoding="utf-8"))


def _validate(doc: dict) -> None:
    validator = jsonschema.Draft202012Validator(_on_disk())
    errors = sorted(validator.iter_errors(doc), key=lambda e: list(e.path))
    assert not errors, "\n".join(
        f"{list(e.path)}: {e.message}" for e in errors[:5]
    )


def test_a_text_document_validates_against_the_published_schema(tmp_path):
    src = tmp_path / "notes.txt"
    src.write_text("First paragraph.\n\nSecond paragraph.\n", encoding="utf-8")
    _validate(_uir_from(tmp_path, src))


def test_an_rtf_document_validates_against_the_published_schema(tmp_path):
    pytest.importorskip("striprtf")
    src = tmp_path / "n.rtf"
    src.write_bytes(rb"{\rtf1\ansi Hello RTF world.}")
    _validate(_uir_from(tmp_path, src))


def test_a_docx_document_validates_against_the_published_schema(tmp_path):
    pytest.importorskip("docx")
    from docx import Document

    src = tmp_path / "r.docx"
    d = Document()
    d.add_heading("Title", level=1)
    d.add_paragraph("Body text.")
    d.save(str(src))
    _validate(_uir_from(tmp_path, src))


def test_a_pptx_document_validates_against_the_published_schema(tmp_path):
    pytest.importorskip("pptx")
    from pptx import Presentation

    src = tmp_path / "d.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Slide Title"
    slide.placeholders[1].text = "Body of the slide."
    prs.save(str(src))
    _validate(_uir_from(tmp_path, src))
