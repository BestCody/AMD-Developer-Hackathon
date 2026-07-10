"""`ingest_any` is the non-PDF ingress; `_extract_pptx_route` is the PPTX walker.

`pipeline.run` used to call `ingest()` for everything, which asserts `%PDF-`
magic bytes -- so every DOCX/PPTX/XLSX died at the door with "is not a PDF"
even though `format_router` classified them correctly and `DocumentConverter`
accepts them natively. Only the PDF and IMAGE routes were ever implemented.
"""
from __future__ import annotations

import zipfile

import pytest

from uir_pipeline.ingest import ingest, ingest_any


def _write_docx(path):
    from docx import Document

    d = Document()
    d.add_heading("Quarterly Report", level=1)
    d.add_paragraph("Revenue grew twelve percent year over year.")
    d.save(str(path))
    return path


def _write_pptx(path, n_slides=2):
    from pptx import Presentation

    prs = Presentation()
    for i in range(n_slides):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"Slide {i + 1} Title"
        slide.placeholders[1].text = f"Body text for slide {i + 1}."
    prs.save(str(path))
    return path


# ---------------------------------------------------------------------------
# ingest_any
# ---------------------------------------------------------------------------

def test_ingest_any_accepts_a_docx(tmp_path):
    pytest.importorskip("docx")
    doc = ingest_any(_write_docx(tmp_path / "report.docx"))
    assert doc.format == "DOCX"
    assert doc.route == "docling"
    assert doc.sha256 and doc.size_bytes > 0
    assert doc.uri.startswith("file://")


def test_ingest_any_routes_pptx_to_the_native_walker(tmp_path):
    pytest.importorskip("pptx")
    doc = ingest_any(_write_pptx(tmp_path / "deck.pptx"))
    assert doc.format == "PPTX"
    assert doc.route == "pptx", "docling returns 0 regions for an unrendered deck"


def test_ingest_any_reports_page_count_zero_not_one(tmp_path):
    """OOXML has no page count until laid out. 0 means 'not applicable'."""
    pytest.importorskip("docx")
    assert ingest_any(_write_docx(tmp_path / "r.docx")).page_count == 0


def test_ingest_any_rejects_an_unsupported_format(tmp_path):
    blob = tmp_path / "mystery.bin"
    blob.write_bytes(b"\x00\x01\x02\x03not a document")
    with pytest.raises(ValueError, match="unsupported format"):
        ingest_any(blob)


def test_ingest_any_rejects_a_missing_file(tmp_path):
    with pytest.raises(FileNotFoundError):
        ingest_any(tmp_path / "nope.docx")


def test_ingest_any_detects_by_content_not_extension(tmp_path):
    """A DOCX renamed to .pdf must still be seen as a DOCX, not trusted blindly."""
    pytest.importorskip("docx")
    src = _write_docx(tmp_path / "real.docx")
    liar = tmp_path / "liar.pdf"
    liar.write_bytes(src.read_bytes())
    assert ingest_any(liar).format == "DOCX"


def test_ingest_still_refuses_a_non_pdf(tmp_path):
    """`ingest` keeps its PDF guarantee; `ingest_any` is a sibling, not a widening."""
    pytest.importorskip("docx")
    with pytest.raises(ValueError, match="is not a PDF"):
        ingest(_write_docx(tmp_path / "r.docx"))


def test_ingest_any_source_metadata_round_trips_to_uir(tmp_path):
    pytest.importorskip("docx")
    doc = ingest_any(_write_docx(tmp_path / "report.docx"))
    source, metadata = doc.to_uir_source_metadata()
    assert source.format == "DOCX"
    assert source.route == "docling"
    assert source.checksum.startswith("sha256:")
    assert metadata.title == "report"  # filename stem, not "(untitled)"


# ---------------------------------------------------------------------------
# _extract_pptx_route
# ---------------------------------------------------------------------------

def test_pptx_route_emits_one_page_per_slide(tmp_path):
    pytest.importorskip("pptx")
    from uir_pipeline.pipeline import _extract_pptx_route

    regions = _extract_pptx_route(_write_pptx(tmp_path / "deck.pptx", n_slides=3))
    assert {r.page for r in regions} == {1, 2, 3}


def test_pptx_route_labels_the_title_placeholder_as_a_heading(tmp_path):
    pytest.importorskip("pptx")
    from uir_pipeline.layout import LayoutLabel
    from uir_pipeline.pipeline import _extract_pptx_route

    regions = _extract_pptx_route(_write_pptx(tmp_path / "deck.pptx", n_slides=1))
    headings = [r for r in regions if r.label is LayoutLabel.HEADING]
    assert len(headings) == 1
    assert headings[0].text == "Slide 1 Title"


def test_pptx_route_emits_validator_safe_bboxes(tmp_path):
    """ChunkNode requires x1 <= x2, y1 <= y2, all within the 0-1000 canvas."""
    pytest.importorskip("pptx")
    from uir_pipeline.pipeline import _extract_pptx_route

    for r in _extract_pptx_route(_write_pptx(tmp_path / "deck.pptx", n_slides=2)):
        x1, y1, x2, y2 = r.bbox
        assert 0 <= x1 <= x2 <= 1000
        assert 0 <= y1 <= y2 <= 1000


def test_pptx_route_skips_empty_shapes(tmp_path):
    pytest.importorskip("pptx")
    from pptx import Presentation

    from uir_pipeline.pipeline import _extract_pptx_route

    path = tmp_path / "sparse.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Only a title"
    slide.placeholders[1].text = "   "  # whitespace only
    prs.save(str(path))

    regions = _extract_pptx_route(path)
    assert len(regions) == 1
    assert all(r.text.strip() for r in regions)


def test_pptx_route_reading_order_is_monotonic(tmp_path):
    pytest.importorskip("pptx")
    from uir_pipeline.pipeline import _extract_pptx_route

    orders = [r.reading_order for r in _extract_pptx_route(_write_pptx(tmp_path / "d.pptx", 3))]
    assert orders == sorted(orders)
    assert len(set(orders)) == len(orders), "reading_order must be unique"


def test_page_texts_from_regions_groups_and_orders_by_page():
    from types import SimpleNamespace

    from uir_pipeline.pipeline import _page_texts_from_regions

    regions = [
        SimpleNamespace(page=2, text="second page"),
        SimpleNamespace(page=1, text="first"),
        SimpleNamespace(page=1, text="also first"),
    ]
    assert _page_texts_from_regions(regions) == [
        (1, "first\nalso first"),
        (2, "second page"),
    ]


def test_xlsx_is_classified_as_the_docling_route(tmp_path):
    """A spreadsheet has no pages at all; docling handles it natively."""
    pytest.importorskip("openpyxl")
    from openpyxl import Workbook

    path = tmp_path / "book.xlsx"
    wb = Workbook()
    wb.active["A1"] = "Region"
    wb.active["B1"] = "Revenue"
    wb.save(str(path))

    doc = ingest_any(path)
    assert doc.format == "XLSX" and doc.route == "docling"
    assert doc.page_count == 0


def test_a_bare_zip_is_not_mistaken_for_ooxml(tmp_path):
    """DOCX/PPTX/XLSX all share ZIP magic; the subtype must be inspected."""
    path = tmp_path / "plain.zip"
    with zipfile.ZipFile(path, "w") as z:
        z.writestr("hello.txt", "not an office document")
    with pytest.raises(ValueError, match="unsupported format"):
        ingest_any(path)


# ---------------------------------------------------------------------------
# The pageless TEXT route
# ---------------------------------------------------------------------------
# format_router sends TXT / MD / CSV / RTF / source code to a "read -> paginate
# -> chunk" lane that never invokes Docling. It was declared (paginate_pageless
# names `pipeline._run_text_route` as a caller) but never written, so these
# files reached `extract_with_docling`, whose allow-list rejects .rtf and .py
# outright: "File format not allowed".

def test_text_route_reads_a_plain_text_file(tmp_path):
    from uir_pipeline.pipeline import _run_text_route

    p = tmp_path / "notes.txt"
    p.write_text("First paragraph.\n\nSecond paragraph.\n", encoding="utf-8")
    regions, pages = _run_text_route(p, "TXT")
    assert [r.text for r in regions] == ["First paragraph.", "Second paragraph."]
    assert pages and pages[0][0] == 1


def test_text_route_handles_source_code(tmp_path):
    """Docling's allow-list rejects .py; the text route must not."""
    from uir_pipeline.pipeline import _run_text_route

    p = tmp_path / "mod.py"
    p.write_text('def f():\n    return 1\n', encoding="utf-8")
    regions, _ = _run_text_route(p, "PY")
    assert regions and "def f()" in regions[0].text


def test_text_route_survives_undecodable_bytes(tmp_path):
    """A stray 0x80 must not fail the document."""
    from uir_pipeline.pipeline import _run_text_route

    p = tmp_path / "weird.txt"
    p.write_bytes(b"before \x80\xff after")
    regions, _ = _run_text_route(p, "TXT")
    assert regions
    assert "before" in regions[0].text and "after" in regions[0].text


def test_text_route_emits_validator_safe_bboxes(tmp_path):
    from uir_pipeline.pipeline import _run_text_route

    p = tmp_path / "n.txt"
    p.write_text("one\n\ntwo\n", encoding="utf-8")
    for r in _run_text_route(p, "TXT")[0]:
        x1, y1, x2, y2 = r.bbox
        assert 0 <= x1 <= x2 <= 1000 and 0 <= y1 <= y2 <= 1000


def test_text_route_reading_order_is_unique_and_monotonic(tmp_path):
    from uir_pipeline.pipeline import _run_text_route

    p = tmp_path / "n.txt"
    p.write_text("a\n\nb\n\nc\n", encoding="utf-8")
    orders = [r.reading_order for r in _run_text_route(p, "TXT")[0]]
    assert orders == sorted(orders) == [1, 2, 3]


def test_text_route_skips_blank_paragraphs(tmp_path):
    from uir_pipeline.pipeline import _run_text_route

    p = tmp_path / "n.txt"
    p.write_text("a\n\n\n\n   \n\nb\n", encoding="utf-8")
    assert [r.text for r in _run_text_route(p, "TXT")[0]] == ["a", "b"]


def test_text_route_decodes_rtf(tmp_path):
    pytest.importorskip("striprtf")
    from uir_pipeline.pipeline import _run_text_route

    p = tmp_path / "n.rtf"
    p.write_bytes(rb"{\rtf1\ansi Hello RTF world.}")
    regions, _ = _run_text_route(p, "RTF")
    assert regions
    text = " ".join(r.text for r in regions)
    assert "Hello RTF world" in text
    assert "\rtf1" not in text, "control words must be stripped"


# ---------------------------------------------------------------------------
# Provenance: Source.route must name the lane that actually ran
# ---------------------------------------------------------------------------

def test_mime_types_do_not_depend_on_the_host_registry(monkeypatch):
    """On Windows, mimetypes reads the registry: .csv -> application/vnd.ms-excel.

    The same file must not carry a different Source.mime_type depending on
    which machine converted it.
    """
    import mimetypes

    from uir_pipeline.ingest import _MIME_BY_FORMAT

    # Poison the stdlib table the way a Windows registry does.
    monkeypatch.setattr(
        mimetypes, "guess_type", lambda *_a, **_k: ("application/vnd.ms-excel", None)
    )
    assert _MIME_BY_FORMAT["CSV"] == "text/csv"
    assert _MIME_BY_FORMAT["RTF"] == "application/rtf"
    assert _MIME_BY_FORMAT["TXT"] == "text/plain"


def test_ingest_any_gives_a_text_file_a_text_mime(tmp_path):
    p = tmp_path / "n.txt"
    p.write_text("hello", encoding="utf-8")
    assert ingest_any(p).mime_type == "text/plain"


def test_ingest_any_gives_csv_a_stable_mime(tmp_path):
    p = tmp_path / "n.csv"
    p.write_text("a,b\n1,2\n", encoding="utf-8")
    assert ingest_any(p).mime_type == "text/csv"


def test_unknown_text_extension_still_gets_a_mime(tmp_path):
    """Code extensions fall through to the stdlib table; any answer beats none."""
    p = tmp_path / "mod.py"
    p.write_text("x = 1\n", encoding="utf-8")
    assert ingest_any(p).mime_type
