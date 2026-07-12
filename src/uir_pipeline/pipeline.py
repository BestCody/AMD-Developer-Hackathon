"""pipeline -- programmatic orchestrator (Phase L).

PLAN.md \u00a79 Phase L exit:
    -- chain ingest -> ocr -> layout -> tables -> chunk -> enrich -> embed -> assemble
    -- provenance block populated with model name, version, and ISO timestamp
    -- emits a single ``UIRV1`` JSON per document
    -- serial processing is fine for MVP

The orchestrator's text path is **IBM Docling** (PLAN OCR follow-up
done): a transformer-based PDF->DoclingDocument emitter that returns
pre-typed sections / tables / figures / math so downstream chunks come
out spatially-aware and column-correct instead of flattened prose.
DOCLING is the only backend -- the previous pdfplumber fast path was
retired because it couldn't preserve column structure. Passing
``UIR_FAST_PATH=pdfplumber`` or ``fast_path="pdfplumber"`` emits a
one-shot deprecation warning and routes through docling.

Weaviate upsert is optional via ``skip_weaviate=True``. When enabled, the
orchestrator (a) ensures both ``UIRChunks_v1`` and ``UIRParentDoc_v1``
collections exist, (b) writes one row per chunk with the prefixed UIR id
stored as a BM25 property, and (c) writes the document-level mean-pool
aggregate to the parent collection.
"""
from __future__ import annotations

import json
import logging
import os
import re
import time
import dataclasses
from dataclasses import dataclass
from datetime import datetime, timezone
from pathlib import Path
from types import SimpleNamespace
from typing import Any

logger = logging.getLogger(__name__)


# ----------------------------------------------------------------------------
# Boilerplate filter
# ----------------------------------------------------------------------------

# Patterns that surface repeatedly on ArXiv-style PDFs because pdfplumber
# concatenates narrow-kerned corporate tokens without explicit spaces.
# Apply post-NER; drop matching entities and any relationships whose
# endpoints reference a dropped entity. Tested against: "Attention Is All
# You Need" (1706.03762) where this filter dropped 21 entities and 60
# relationships from the 411/736 baseline (Tier 3 fix #2 follow-up).
_BOILERPLATE_RE: tuple[re.Pattern[str], ...] = (
    # The ArXiv permission block ("Google hereby grants permission to..."),
    # including the token-stripped variant pdfplumber emits. The trailing
    # "to" is suffixed via `(?:to)?` not a `\b` so the variant
    # ``"Googleherebygrantspermissionto"`` (pdfplumber concatenates without
    # spaces) still matches: ``\b`` between two word-class chars (``n`` /
    # ``t``) would not fire, so a naïve ``\b…\b`` pattern misses the case.
    re.compile(r"\bgoogle\s+hereby\s+grants?\s+permissions?\b", re.IGNORECASE),
    re.compile(r"\bgoogleherebygrantspermission(?:to)?\b", re.IGNORECASE),
    # Google Research / Brain affiliations repeated in headers/footers.
    re.compile(r"\bgoogle\s*brain\b", re.IGNORECASE),
    re.compile(r"\bgoogle\s*research\b", re.IGNORECASE),
    # Mountain View / corporate addr fragments.
    re.compile(r"\bmountain\s+view\b", re.IGNORECASE),
    # Standard copyright token clusters.
    re.compile(r"\bcopyright\s+\(c\)\s*\d{4}\b", re.IGNORECASE),
)


def _is_boilerplate(text: str) -> bool:
    """Return True iff ``text`` matches any :data:`_BOILERPLATE_RE` pattern.

    Used by :func:`run` post-enrichment to drop noisy generic entities that
    pdfplumber surfaces from ArXiv corporate footers/permission headers.
    """
    return any(p.search(text) for p in _BOILERPLATE_RE)


# ----------------------------------------------------------------------------
# Public result
# ----------------------------------------------------------------------------

@dataclass(frozen=True)
class PipelineResult:
    """Per-document pipeline outcome.

    ``umr_path`` is the Universal Markdown Representation companion
    file ``{uir_id}.umr.md`` emitted alongside ``{uir_id}.uir.json``.
    Always populated (Phase 17 §UMR); the agent-facing view of the
    document. ``entity_count`` reflects the post-boilerplate-filter
    entity count when ``include_semantics`` is True (default False),
    else 0 to surface the empty semantics block honestly. The new
    fields are defaulted so test stubs / older callers that don't
    kwarg-pass them keep working.
    """
    uir_id: str
    out_path: Path
    umr_path: Path | None = None
    chunk_count: int = 0
    entity_count: int = 0
    elapsed_seconds: float = 0.0


# ----------------------------------------------------------------------------
# Page-text extraction (fast path: pdfplumber; real OCR is a one-line swap)
# ----------------------------------------------------------------------------

# _get_page_text removed -- page text now comes from DoclingResult.page_texts

# ----------------------------------------------------------------------------
# Fast-path resolution + Docling shims (PLAN §17 §OCR follow-up)
# ----------------------------------------------------------------------------

def _resolve_fast_path(fast_path: str | None) -> str:
    """Resolve the active fast-path backend. Always returns ``"docling"``.

    Legacy callers that pass ``UIR_FAST_PATH=pdfplumber`` or
    ``fast_path="pdfplumber"`` are redirected to docling with a one-shot
    deprecation warning. The previous pdfplumber fast path was a
    column-naive text scraper; the docling layout-model backend has fully
    replaced it (PLAN OCR follow-up done).

    Priority: explicit ``fast_path`` arg > ``UIR_FAST_PATH`` env var >
    ``"docling"`` (production default). Unknown values log a warning and
    default to docling.
    """
    raw = (
        (fast_path or os.environ.get("UIR_FAST_PATH", "") or "docling")
        .strip()
        .lower()
    )
    if raw == "pdfplumber":
        logger.warning(
            "UIR_FAST_PATH=%r is deprecated -- routing to docling (pdfplumber fast path removed)",
            raw,
        )
        return "docling"
    if raw != "docling":
        logger.warning("unknown UIR_FAST_PATH=%r -- defaulting to docling", raw)
        return "docling"
    return raw


class ImageAnalysisError(RuntimeError):
    """The IMAGE route could not describe the image.

    `run_image_pipeline` reports failure in `ImagePipelineResult.error` rather
    than raising, so the orchestrator has to translate. Without this the caller
    saw a well-formed result whose `out_path` pointed at a file that was never
    written.
    """


class AudioAnalysisError(RuntimeError):
    """The AUDIO route could not transcribe the audio file.

    `run_audio_pipeline` reports failure in `AudioPipelineResult.error` rather
    than raising, so the orchestrator has to translate. Mirrors the same
    pattern as `ImageAnalysisError`.
    """


class VideoAnalysisError(RuntimeError):
    """The VIDEO route could not process the video file.

    `run_video_pipeline` reports failure in `VideoPipelineResult.error` rather
    than raising, so the orchestrator has to translate. Mirrors the same
    pattern as `ImageAnalysisError` and `AudioAnalysisError`.
    """


class _NoFigureSource(Exception):
    """The chosen route produced no figure regions to caption.

    Not an error: PPTX_NATIVE never runs Docling, so there is no
    ``DoclingResult`` carrying figure bboxes. Signalled as an exception only
    so the caption stage's existing fail-soft structure can skip it without
    logging a failure.
    """


def _page_texts_from_regions(regions: list[Any]) -> list[tuple[int, str]]:
    """Collapse regions into ``[(page, joined_text), ...]``, page-ordered.

    The Docling route gets this from ``dr.page_texts``; routes that build
    regions themselves still owe the chunker the same page-level view.
    """
    by_page: dict[int, list[str]] = {}
    for r in regions:
        by_page.setdefault(int(r.page), []).append(r.text)
    return [(page, "\n".join(texts)) for page, texts in sorted(by_page.items())]


def _read_text_file(path: Path) -> str:
    """Read a text file without ever raising on encoding.

    A stray 0x80 byte in a source file must not fail the whole document, so
    undecodable bytes become U+FFFD rather than a UnicodeDecodeError. UTF-8
    first because everything modern is; latin-1 never fails but silently
    mojibakes, so it is not a fallback worth having.
    """
    return path.read_text(encoding="utf-8", errors="replace")


def _split_paragraphs(text: str) -> list[str]:
    """Split on blank lines, but never inside a ``` fenced block.

    A markdown or notebook code block routinely contains blank lines. A naive
    ``re.split(r"\\n\\s*\\n", ...)`` tears it in half, leaving an unterminated
    fence in one chunk and an orphan ``` in the next.
    """
    paragraphs: list[str] = []
    current: list[str] = []
    in_fence = False

    def _flush() -> None:
        block = "\n".join(current).strip()
        if block:
            paragraphs.append(block)
        current.clear()

    for line in text.splitlines():
        if line.lstrip().startswith("```"):
            in_fence = not in_fence
            current.append(line)
            if not in_fence:  # closing fence ends the block
                _flush()
            continue
        if not in_fence and not line.strip():
            _flush()
            continue
        current.append(line)
    _flush()
    return paragraphs


def _notebook_to_text(path: Path) -> str:
    """Flatten a Jupyter notebook into prose + fenced code.

    Docling's allow-list carries no notebook format, so `.ipynb` used to fail
    with "File format not allowed". It is JSON, but dumping the JSON would bury
    the prose under `"cell_type"` / `"metadata"` noise and embed base64 image
    outputs in the embeddings.

    Markdown cells pass through as-is. Code cells become fenced blocks so the
    chunker keeps them whole. Outputs are dropped: they are often megabytes of
    base64 PNG, and an execution result is not part of the document's meaning.
    """
    try:
        nb = json.loads(_read_text_file(path))
    except json.JSONDecodeError as exc:
        raise ValueError(f"{path.name} is not valid notebook JSON: {exc}") from exc

    lang = (
        (nb.get("metadata") or {}).get("kernelspec", {}).get("language")
        or "python"
    )
    blocks: list[str] = []
    for cell in nb.get("cells") or []:
        source = cell.get("source") or []
        # nbformat stores source as a list of lines *or* a single string.
        text = ("".join(source) if isinstance(source, list) else str(source)).strip()
        if not text:
            continue
        kind = cell.get("cell_type")
        if kind == "markdown":
            blocks.append(text)
        elif kind == "code":
            blocks.append(f"```{lang}\n{text}\n```")
        # raw cells carry no rendering semantics; keep their text verbatim.
        elif kind == "raw":
            blocks.append(text)
    return "\n\n".join(blocks)


def _run_text_route(path: Path, fmt: str) -> tuple[list[Any], list[tuple[int, str]]]:
    """The pageless route: read -> paginate -> regions. No Docling.

    ``format_router`` sends TXT / MD / CSV / RTF / source code here. Docling
    cannot open most of them -- `.rtf` and `.py` fail its allow-list outright
    ("File format not allowed") -- and they have no page geometry to recover,
    so a layout model has nothing to contribute.

    Pages are synthesized by :func:`chunk.paginate_pageless` at ~2000 BGE
    tokens each, matching the PDF contract of 1-based page numbers. Within a
    page, blank-line-separated paragraphs become regions so the chunker sees
    real paragraph boundaries instead of one undifferentiated blob.

    Bounding boxes are the full canvas: a text file has no geometry, and
    claiming a narrower box would be inventing provenance.
    """
    from uir_pipeline.chunk import paginate_pageless
    from uir_pipeline.layout import LayoutLabel, LayoutRegion

    fmt_upper = fmt.upper()
    if fmt_upper == "RTF":
        # striprtf decodes the control words; it also paginates for us.
        from uir_pipeline.ingest_rtf import ingest_rtf

        _doc, page_pairs = ingest_rtf(path)
    elif fmt_upper == "IPYNB":
        page_pairs = paginate_pageless(_notebook_to_text(path))
    else:
        page_pairs = paginate_pageless(_read_text_file(path))

    # Markdown (and notebook markdown cells) declare headings with `#`. Docling
    # used to label them for us; on this lane nothing would, and the chunker's
    # section-path tracking only recognises numbered headings.
    markdownish = fmt_upper in ("MD", "MARKDOWN", "IPYNB")

    regions: list[Any] = []
    order = 0
    for page_no, page_text in page_pairs:
        for para in _split_paragraphs(page_text):
            order += 1
            is_heading = (
                markdownish
                and para.startswith("#")
                and not para.startswith("#!")  # a shebang is not a heading
                and "\n" not in para           # a heading is one line
            )
            regions.append(LayoutRegion(
                label=LayoutLabel.HEADING if is_heading else LayoutLabel.PARAGRAPH,
                text=para,
                # Read verbatim off disk; nothing was inferred, so nothing is
                # uncertain.
                confidence=1.0,
                bbox=(0, 0, 1000, 1000),
                page=int(page_no),
                reading_order=order,
            ))
    return regions, list(page_pairs)


def _extract_pptx_route(path: Path) -> list[Any]:
    """Walk a ``.pptx`` with python-pptx and return :class:`LayoutRegion` list.

    Referenced by ``format_router.classify_route``, which sends PPTX here
    rather than to Docling: Docling's layout model runs on rendered page
    images, and a python-pptx-generated deck has no rendering, so it returns
    zero regions. A slide is already structured -- title placeholder, then
    body placeholders -- so a native walk beats a layout model that has
    nothing to look at.

    One slide == one page. Shapes carry EMU offsets, not a 0-1000 canvas,
    so bboxes are scaled against the presentation's own slide dimensions
    rather than assumed to be letter-sized.
    """
    from pptx import Presentation
    from pptx.util import Emu

    from uir_pipeline.layout import LayoutLabel, LayoutRegion

    prs = Presentation(str(path))
    slide_w = int(prs.slide_width or Emu(9144000))
    slide_h = int(prs.slide_height or Emu(6858000))

    def _canvas(shape: Any) -> tuple[int, int, int, int]:
        try:
            left, top = int(shape.left or 0), int(shape.top or 0)
            width, height = int(shape.width or 0), int(shape.height or 0)
        except (TypeError, ValueError):
            return (0, 0, 0, 0)
        x1 = max(0, min(1000, round(left * 1000 / slide_w)))
        y1 = max(0, min(1000, round(top * 1000 / slide_h)))
        x2 = max(0, min(1000, round((left + width) * 1000 / slide_w)))
        y2 = max(0, min(1000, round((top + height) * 1000 / slide_h)))
        # ChunkNode's validator requires x1 <= x2 and y1 <= y2.
        if x1 > x2:
            x1, x2 = x2, x1
        if y1 > y2:
            y1, y2 = y2, y1
        return (x1, y1, x2, y2)

    regions: list[Any] = []
    order = 0
    for page_no, slide in enumerate(prs.slides, start=1):
        # `slide.shapes.title` builds a fresh proxy on every access, so
        # `shape is slide.shapes.title` is always False. `shape_id` is the
        # stable identity.
        title_id = None
        try:
            title = slide.shapes.title
            title_id = title.shape_id if title is not None else None
        except (AttributeError, ValueError):
            pass
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            text = (shape.text_frame.text or "").strip()
            if not text:
                continue
            order += 1
            is_title = title_id is not None and shape.shape_id == title_id
            regions.append(LayoutRegion(
                label=LayoutLabel.HEADING if is_title else LayoutLabel.PARAGRAPH,
                text=text,
                # A placeholder's role is declared in the file, not inferred
                # by a model, so there is nothing to be uncertain about.
                confidence=1.0,
                bbox=_canvas(shape),
                page=page_no,
                reading_order=order,
            ))
    return regions


def _docling_to_table_draft(t: dict[str, Any]) -> Any:
    """Synthesize a :class:`TableDraft` from a docling tables[].dict.

    Docling natively exports tables as GitHub-flavored markdown with
    a ``|``-separated header row preserved. :attr:`TableDraft.row_count`
    and :attr:`TableDraft.col_count` are derived from the rendered
    markdown: pipes_count-1 on the first non-separator row for cols;
    count of non-separator ``|``-led lines for rows. The
    :attr:`TableDraft.confidence` field is a static ``0.9`` because the
    standard ``DocumentConverter`` output doesn't expose a per-region
    logprob -- downstream consumers can apply their own threshold.
    """
    from uir_pipeline.tables import TableDraft
    md = t["markdown"]
    rows = [
        r for r in md.splitlines()
        if r.strip().startswith("|") and "---" not in r
    ]
    col_count = max(0, rows[0].count("|") - 1) if rows else 0
    row_count = len(rows)
    return TableDraft(
        page_number=int(t["page"]),
        bbox=tuple(t["bbox"]),
        markdown=md,
        row_count=row_count,
        col_count=col_count,
        confidence=0.9,
    )


def _is_weaviate_unavailable(exc: BaseException) -> bool:
    """True iff ``exc`` means "no Weaviate server here", not "our code is wrong".

    Used to decide whether the optional upsert stage may fail soft. Only
    transport/startup conditions qualify. A schema mismatch or a rejected
    object is a defect and must propagate.
    """
    if isinstance(exc, ImportError):  # weaviate-client not installed
        return True
    try:
        from weaviate.exceptions import (
            WeaviateConnectionError,
            WeaviateGRPCUnavailableError,
            WeaviateStartUpError,
            WeaviateTimeoutError,
        )
    except ImportError:  # pragma: no cover -- client absent entirely
        return True
    return isinstance(exc, (
        WeaviateConnectionError,
        WeaviateGRPCUnavailableError,
        WeaviateStartUpError,
        WeaviateTimeoutError,
    ))


# ----------------------------------------------------------------------------
# Public API
# ----------------------------------------------------------------------------

def run(
    input_path: str | Path,
    output_dir: str | Path,
    *,
    skip_weaviate: bool = False,
    dry_run: bool = False,
    with_embeddings: bool = True,
    page_numbers: list[int] | None = None,
    on_progress: Any | None = None,
    include_semantics: bool = False,
    fast_path: str | None = None,
    intent: str | None = None,
) -> PipelineResult:
    """Drive the full pipeline on one PDF and return a :class:`PipelineResult`.

    Parameters:
        input_path: PDF file path.
        output_dir: Where to write ``{uir_id}.uir.json`` (and the companion
            ``{uir_id}.umr.md`` -- always emitted alongside).
        skip_weaviate: True -> don't upsert to Weaviate (default: false).
        dry_run: True -> don't write JSON or Weaviate (default: false).
        with_embeddings: True -> compute BGE embeddings (default: true).
            False -> skip the embed step (faster, useful for tests).
        page_numbers: 1-based list of pages to process (``None`` = all).
        on_progress: optional callback ``fn(stage: str, percent: int)``.
        include_semantics: True -> emit the verbose ``semantics`` block in
            the UIR JSON (entities + relationships + topics; can be 400+
            + 700+ lines on a real arXiv doc). Default: ``False`` --
            semantics is OMITTED from the JSON output so agent-facing
            consumers receive a clean payload.
        fast_path: Per-page text-extraction backend. ``"docling"`` is the
            only legal value; ``"pdfplumber"`` is accepted as a
            deprecated alias and emits a one-shot warning.
        intent: Optional intent/query for image files. When the input is
            an IMAGE and intent is provided, it's passed to the Fireworks
            AI vision model instead of generating a generic description.
            Ignored for non-image routes.

    Returns:
        A :class:`PipelineResult` with :attr:`PipelineResult.umr_path`
        populated (the companion ``.umr.md`` file is always written
        alongside ``.uir.json`` so agents have the same view that
        :file:`templates/console.html` surfaces).
    """
    t0 = time.monotonic()
    p = Path(input_path)
    output_dir = Path(output_dir)

    # Route check: IMAGE files go through the Fireworks AI vision pipeline
    # (separate from the PDF/DOCX/TEXT pipeline below).
    from uir_pipeline.format_router import FormatRoute, route as _format_route

    fmt, froute = _format_route(p)
    if froute is FormatRoute.IMAGE:
        from uir_pipeline.image_pipeline import run_image_pipeline

        img_result = run_image_pipeline(
            p,
            output_dir=output_dir,
            intent=intent,
            dry_run=dry_run,
            on_progress=on_progress,
        )
        if img_result.error:
            # `run_image_pipeline` reports failure in a field, not an
            # exception. Folding that into `chunk_count=0` and returning a
            # success-shaped result made the CLI log
            # "done chart.png: chunks=0 -> <path>" and exit 0 while writing no
            # file at all, and made the web job report `done` with a `result`
            # that 404s. Raise instead: the caller already turns this into a
            # failed job (web) or a non-zero exit (CLI).
            raise ImageAnalysisError(img_result.error)
        # Synthesise a PipelineResult from the ImagePipelineResult so the
        # calling CLI/web layer receives the same shape it expects.
        return PipelineResult(
            uir_id=img_result.uir_id,
            out_path=img_result.out_path,
            umr_path=img_result.umr_path,
            chunk_count=1,
            entity_count=0,
            elapsed_seconds=img_result.elapsed_seconds,
        )

    if froute is FormatRoute.AUDIO:
        from uir_pipeline.audio_pipeline import run_audio_pipeline

        audio_result = run_audio_pipeline(
            p,
            output_dir=output_dir,
            dry_run=dry_run,
            on_progress=on_progress,
        )
        if audio_result.error:
            raise AudioAnalysisError(audio_result.error)
        return PipelineResult(
            uir_id=audio_result.uir_id,
            out_path=audio_result.out_path,
            umr_path=audio_result.umr_path,
            chunk_count=audio_result.chunk_count,
            entity_count=audio_result.entity_count,
            elapsed_seconds=audio_result.elapsed_seconds,
        )

    if froute is FormatRoute.VIDEO:
        from uir_pipeline.video_pipeline import run_video_pipeline

        video_result = run_video_pipeline(
            p,
            output_dir=output_dir,
            dry_run=dry_run,
            on_progress=on_progress,
        )
        if video_result.error:
            raise VideoAnalysisError(video_result.error)
        return PipelineResult(
            uir_id=video_result.uir_id,
            out_path=video_result.out_path,
            umr_path=video_result.umr_path,
            chunk_count=video_result.chunk_count,
            entity_count=video_result.entity_count,
            elapsed_seconds=video_result.elapsed_seconds,
        )

    from uir_pipeline.chunk import chunk_text
    from uir_pipeline.embed import (
        derive_doc_id,
        embed_texts,
        ensure_collections,
        mean_pool_vectors,
        upsert_chunks,
        upsert_parent_doc,
    )
    from uir_pipeline.enrich import EnrichmentResult, enrich_chunks
    from uir_pipeline.ingest import DocumentInput, ingest, ingest_any
    from uir_pipeline.logging_config import (
        attach_doc_log,
        configure,
        detach_doc_log,
    )
    from uir_pipeline.uir_schema import (
        ChunkNode,
        Entity,
        ExtractionProvenance,
        NormalizationProvenance,
        Provenance,
        Relationship,
        Semantics,
        Structure,
        StructureNode,
        UIRV1,
    )
    from uir_pipeline.utils import (
        deterministic_node_id,
        strip_uir_prefix,
    )

    configure()

    # Stage 1: ingest
    def _progress(stage: str, pct: int, **meta: Any) -> None:
        """Emit a stage-progress event.

        ``meta`` is forwarded to ``on_progress(stage, pct, **meta)`` so
        downstream consumers (the web runner, the LAN server, integration
        tests) can surface per-stage diagnostics like
        ``caption_records_empty=3`` without inheriting the orchestrator's
        internal state. Logged alongside the ``(pct)`` field.
        """
        logger.info("pipeline.stage %s (%d%%) meta=%s", stage, pct, meta)
        if on_progress is not None:
            try:
                on_progress(stage, pct, **meta)
            except Exception:
                pass
    _progress("ingest", 5)
    # `ingest` is the PDF ingress: it asserts %PDF- magic and reads page count
    # with pypdf. Neither applies to OOXML, so non-PDF routes take the generic
    # ingress. `froute` was already resolved above for the IMAGE check.
    doc: DocumentInput = ingest(p) if froute is FormatRoute.PDF else ingest_any(p)
    doc_id = derive_doc_id(doc.uri)
    log_dir = output_dir.parent / "logs" if output_dir.name != "logs" else output_dir
    log_handler = attach_doc_log(strip_uir_prefix(doc_id), log_dir)
    try:
        logger.info("ingested %s: %d pages, sha256=%s", doc.uri, doc.page_count, doc.sha256[:12])

        # Stage 2/3/4/5 -- fast_path routing (PLAN §17 §OCR follow-up).
        # Production default is ``docling``: IBM Docling emits pre-typed
        # sections / tables / figures / math natively so chunks come out
        # structured instead of flattened prose. The heuristic
        # :class:`LayoutClassifier` is skipped on that branch.
        # :file:`tests/conftest.py` pins ``UIR_FAST_PATH=pdfplumber`` so
        # CI doesn't pay the 2 GB HuggingFace weight download; pass
        # ``fast_path="docling"`` (or set the env var) on a real dev
        # machine to opt in. When the docling branch raises
        # :class:`DoclingUnavailable` (missing dep OR HF model load
        # failure), the orchestrator cascades to pdfplumber so legacy
        # fixtures keep emitting valid UIR.
        _resolve_fast_path(fast_path)  # side-effect: deprecation warning
        all_regions: list = []
        table_drafts: list = []
        page_text_pairs: list[tuple[int, str]] = []
        # _resolve_fast_path is always "docling" now (pdfplumber is
        # rejected with a deprecation warning). Stage the docling call
        # early so the same DoclingResult feeds both chunk assembly
        # (all_regions / tables / page_texts) AND the caption stage
        # (figures with bboxes from dr.pictures).
        # Docling is the only backend. ``DoclingUnavailable`` propagates
        # naturally here -- the previous try/except that re-raised was a
        # literal no-op once the pdfplumber cascade was deleted. No
        # silent fall-back to a column-naive path; operator sees the
        # real cause (missing dep OR HF model load failure).
        _progress("docling_extract", 18)
        from uir_pipeline.docling_extract import (
            extract_with_docling,
        )
        from uir_pipeline.layout import LayoutLabel, LayoutRegion

        docling_result = None
        if froute is FormatRoute.PPTX_NATIVE:
            # Docling's layout model reads rendered page images; a
            # python-pptx deck has no rendering, so it returns 0 regions.
            all_regions = _extract_pptx_route(p)
            page_text_pairs = _page_texts_from_regions(all_regions)
        elif froute is FormatRoute.TEXT:
            # Pageless: docling's allow-list rejects .rtf and source code
            # outright, and plain text has no layout to recover.
            all_regions, page_text_pairs = _run_text_route(p, doc.format)
        else:
            # PDF and the DOCLING route (DOCX / XLSX / EPUB / HTML / ...):
            # DocumentConverter accepts all of them natively.
            dr = extract_with_docling(p)
            docling_result = dr  # forwarded to figure-caption stage
            all_regions = [
                LayoutRegion(
                    label=LayoutLabel(r["label"]),
                    text=r["text"],
                    confidence=0.9,
                    bbox=tuple(r["bbox"]),
                    page=int(r["page"]),
                    reading_order=i + 1,
                )
                for i, r in enumerate(dr.regions)
            ]
            table_drafts = [_docling_to_table_draft(t) for t in dr.tables]
            page_text_pairs = list(dr.page_texts)

        # Record the lane that actually ran. `ingest_any` copied the router's
        # classification, which is a *request*, not a result -- a .txt was
        # landing in the UIR as `route="text"` while docling did the work.
        # `page_count` is 0 out of `ingest_any` (OOXML has none until laid
        # out); the pageless routes synthesize pages, so report those.
        _actual_route = (
            "pdf" if froute is FormatRoute.PDF
            else "pptx" if froute is FormatRoute.PPTX_NATIVE
            else "text" if froute is FormatRoute.TEXT
            else "docling"
        )
        if doc.route != _actual_route or (doc.page_count == 0 and page_text_pairs):
            doc = dataclasses.replace(
                doc,
                route=_actual_route,
                page_count=doc.page_count or len(page_text_pairs),
            )

        _progress(
            "layout", 45,
            fast_path=_actual_route,
            region_count=len(all_regions),
            table_count=len(table_drafts),
        )

        # Stage 5.5 (Tier 3, per PLAN_TIER3.md): image captioning.
        # Detects figure regions via pdfplumber.Page.images, renders each
        # crop with PyMuPDF, runs Florence-2-base for a structured caption,
        # and emits ChunkNode-compatible shims that share the BGE-embedding
        # pipeline with text chunks. Fail-soft: any exception here is logged
        # and the document emits without figure captions (no UIR schema break).
        _progress("figure_caption", 60)
        figure_chunk_shims: list[Any] = []
        # Counter exposed via on_progress so the UI / integration tests
        # can surface "X figures, Y captioned, Z empty" -- prevents silent
        # loss when Florence-2 fail-softs on bad crops (Tier 3 fix #4).
        caption_records_total = 0
        caption_records_with_text = 0
        try:
            if docling_result is None:
                # PPTX_NATIVE never runs Docling, so there are no figure
                # bboxes to caption. Skip rather than convert a second time.
                raise _NoFigureSource
            from uir_pipeline.caption import caption_figures_in_pdf
            from uir_pipeline.utils import count_tokens as _bpe_count_tokens
            # Forward the same DoclingResult so caption uses figure bboxes
            # from dr.pictures (no pdfplumber path) -- and the source path so
            # it can render the crops. Without `pdf_path` the renderer is
            # skipped entirely and every `image_b64` comes back None; passing
            # both reuses the conversion rather than re-running Docling.
            #
            # Crop rendering goes through PyMuPDF, which opens PDFs. On the
            # DOCLING route (DOCX/XLSX) the bboxes are real but the source is
            # not a PDF, so pass no path: captions still come back, `image_b64`
            # is None, and PyMuPDF is never handed a file it cannot open.
            figure_records = caption_figures_in_pdf(
                p if froute is FormatRoute.PDF else None,
                docling_result=docling_result,
                page_numbers=page_numbers,
            )
            for fig in (figure_records or []):
                caption_records_total += 1
                cap = (fig.get("caption") or "").strip()
                if not cap:
                    continue
                caption_records_with_text += 1
                figure_chunk_shims.append(SimpleNamespace(
                    text=cap,
                    # BPE token count, not word-split: Florence-2's <MORE_DETAILED_CAPTION>
                    # output is subword-tokenized downstream (BGE embedder), so
                    # keeping the figure ChunkNode on the same scale avoids a
                    # UI badge mismatch and keeps BGE chunk-overlap stitching coherent.
                    token_count=_bpe_count_tokens(cap),
                    page=int(fig["page"]),
                    bbox=tuple(fig["bbox_canvas"]),
                    # 0.8 heuristic floor: Florence-2 doesn't expose per-image
                    # logprob at its API surface, and CaptionerBeam scores aren't
                    # directly translatable to a confidence scalar. PLAN_TIER3
                    # risk 8 flagged confidence propagation as deferred. Re-tune
                    # once we collect a labelled figure-caption dataset.
                    confidence=0.8,
                    modal_features={
                        # Tier 1.5 #2: canonicalize the label against
                        # LayoutLabel. ``caption`` is a first-class label
                        # (the chunk carries the caption TEXT -- ``figure``
                        # is reserved for raw figure regions without text).
                        "intent": {"region_kind": "caption"},
                        "figure": {
                            "image_b64": fig.get("image_b64"),
                            "caption_prompt": fig.get("caption_prompt"),
                            "caption_model": fig.get("caption_model"),
                        },
                    },
                ))
        except _NoFigureSource:
            # Not a failure: this route produced no figure bboxes to caption.
            # Distinct from the handler below so it never logs as an error.
            _progress(
                "figure_caption", 60,
                caption_records_total=0, caption_records_with_text=0,
                caption_records_empty=0, skipped="route has no figure source",
            )
        except Exception as exc:
            logger.warning("figure caption stage failed (fail-soft): %s", exc)
            _progress(
                "figure_caption", 60,
                caption_records_total=0, caption_records_with_text=0,
                caption_records_empty=0, error=str(exc),
            )
        else:
            # Emit end-of-stage progress with the actual counts so the
            # caller (web UI / integration test) can diagnose silent loss.
            _progress(
                "figure_caption", 60,
                caption_records_total=caption_records_total,
                caption_records_with_text=caption_records_with_text,
                caption_records_empty=caption_records_total - caption_records_with_text,
            )

        # Stage 6: chunking -- union of layout regions + table markdown + figure captions.
        # Tier 1 intent metadata: walk regions top-down, track ``section_path``
        # state by detecting numbered headings (e.g. ``"3.2 Multi-Head Attention"``),
        # and attach ``region_kind`` (= LayoutLabel.value) to each emitted chunk.
        # The heading regex is anchored at the start of the line and conservative
        # -- only structural headings (numeric prefix; unnumbered items like
        # ``Abstract`` / ``References`` / ``Acknowledgments`` take their literal
        # text as the path) trigger an update. Cross-chunk linking
        # (``preceding_chunk_id`` / ``following_chunk_id``) is wired in
        # Stage 9 below once deterministic chunk IDs are assigned.
        _progress("chunk", 70)
        all_chunks: list[Any] = []
        # Match structural numbering: ``3``, ``3.2``, ``3.2.1`` followed by
        # ``.`` or whitespace, then the section title. Conservative --
        # we never update ``_section_path`` on a non-match, so prose that
        # happens to start with a year ("2024 saw...") won't trigger.
        _section_heading_re = re.compile(
            r"^\s*(\d+(?:\.\d+)*)[\.\s]+(\S.{2,})$"
        )
        _section_path = ""  # starts empty; first heading lights it up
        n_regions = len(all_regions)
        for i, region in enumerate(all_regions):
            label_str = (
                region.label.value if hasattr(region.label, "value")
                else str(region.label)
            )
            if label_str == "heading":
                m = _section_heading_re.match(region.text.strip())
                if m:
                    _section_path = m.group(1)
                else:
                    # Unnumbered heading (Abstract, References, etc.) -- use
                    # the literal text as the path so it's still queryable
                    # via intent-shaped queries.
                    _section_path = region.text.strip().rstrip(".").strip()
            # A region is the LAST of its section if (a) the next region is
            # a new heading OR (b) this is the last region in the document.
            next_label = (
                str(all_regions[i + 1].label.value)
                if i + 1 < n_regions and hasattr(all_regions[i + 1].label, "value")
                else ""
            )
            is_last_of_section = (
                i == n_regions - 1 or next_label == "heading"
            )
            all_chunks.extend(chunk_text(
                region.text,
                page=region.page,
                bbox=region.bbox,
                region_kind=label_str,
                section_path=(_section_path or None),
                is_section_first=(label_str == "heading"),
                is_section_last=bool(is_last_of_section),
            ))
        for table in table_drafts:
            all_chunks.extend(chunk_text(
                table.markdown,
                page=table.page_number,
                bbox=table.bbox,
                region_kind="table",
            ))
        all_chunks.extend(figure_chunk_shims)  # Tier 3 captions get BGE vectors
        # Drop residual 1-3 char noise chunks (Fix Plan item #4). These
        # are axis-tick fragments such as ``"0"`` / ``"0.5"`` / ``"##"``
        # that leak through ``LayoutClassifier`` from figure regions. They
        # destroy sentence-level retrieval signal without adding context.
        # Real arxiv sentences are >= 4 chars in length, so 4 is a safe
        # floor (we keep "BERT"-class acronym chunks, but drop "0.5").
        all_chunks = [ck for ck in all_chunks if len(ck.text.strip()) >= 4]
        if not all_chunks and page_text_pairs:
            # No regions / no tables -- chunk the whole document text.
            full_text = " ".join(text for _, text in page_text_pairs if text)
            all_chunks = chunk_text(full_text, page=1)

        # Stage 7: enrich (NER + co-occurrence)
        _progress("enrich", 80)
        enrichment = enrich_chunks([c.text for c in all_chunks])
        # Filter Arxiv-style boilerplate entities that pdfplumber emits
        # (concatenated narrow-kerned tokens from the permission footer
        # and corporate affiliation blocks). Drop the matching entities
        # and any relationships whose endpoints reference them. The
        # pattern set covers the documented ArXiv noise -- extend here if a
        # new boilerplate source shows up in the entity-quality log.
        before_entity_count = len(enrichment.entities)
        kept_entities = [e for e in enrichment.entities if not _is_boilerplate(e.text)]
        kept_text = {e.text for e in kept_entities}
        kept_relations = [
            r for r in enrichment.relationships
            if r.from_text in kept_text and r.to_text in kept_text
        ]
        dropped_entities = before_entity_count - len(kept_entities)
        dropped_relations = len(enrichment.relationships) - len(kept_relations)
        if dropped_entities or dropped_relations:
            logger.info(
                "boilerplate-filter: dropped %d entities and %d relationships (arXiv footer noise)",
                dropped_entities, dropped_relations,
            )
            _progress(
                "enrich", 80,
                dropped_entities=dropped_entities,
                dropped_relations=dropped_relations,
                entity_count=len(kept_entities),
                relationship_count=len(kept_relations),
            )
        # EnrichmentResult is a frozen dataclass -- rebuild instead of
        # mutating so downstream consumers see the post-filter view in one
        # place. Topics carry through unmodified.
        enrichment = EnrichmentResult(
            entities=kept_entities,
            relationships=kept_relations,
            topics=enrichment.topics,
        )

        # Stage 8: embed (BGE-small 384-d)
        _progress("embed", 90)
        if with_embeddings and all_chunks:
            try:
                vectors = embed_texts([c.text for c in all_chunks])
            except Exception as exc:
                logger.warning("embed failed (%s) -- writing chunks without vectors", exc)
                vectors = None
        else:
            vectors = None

        # Stage 9: assemble UIRV1
        _progress("assemble", 95)
        source, metadata = doc.to_uir_source_metadata()
        # Override the page_count to match what ingest saw.
        metadata = metadata.model_copy(update={"page_count": doc.page_count})

        # Build chunk nodes
        chunk_nodes: list[ChunkNode] = []
        chunk_ids: list[str] = []
        for i, ck in enumerate(all_chunks):
            ck_id = deterministic_node_id("chunk", doc_id, i, ck.text[:64])
            chunk_ids.append(ck_id)
            modal_features = dict(ck.modal_features) if ck.modal_features else {}
            if vectors is not None and i < len(vectors.vectors):
                # Persist the actual float vector onto the chunk so
                # intent_filter can rank by cosine similarity without
                # re-loading BGE on every intent call. Rounded to 6 dp = the
                # cosine ranking is unaffected (order is preserved; cosine
                # range stays inside [-1, 1]). Older UIRs (pre-fix-#1) will
                # simply lack this key -- intent_filter falls back to
                # keyword-only match gracefully.
                modal_features["vector"] = {
                    "dim": vectors.dim,
                    "model": "BAAI/bge-small-en-v1.5",
                    "chunk_index": i,
                    "embedding": [round(float(v), 6) for v in vectors.vectors[i]],
                }
            chunk_nodes.append(ChunkNode(
                id=ck_id,
                type="chunk",
                text=ck.text,
                token_count=ck.token_count,
                page=ck.page,
                bounding_box=ck.bbox,
                confidence=ck.confidence,
                modal_features=modal_features,
            ))
        # Tier 1.C: wire consecutive ``preceding_chunk_id`` and
        # ``following_chunk_id`` per chunk. We use consecutive wiring (not
        # cross-section jumps) so the co-occurrence sliding window in the
        # enrich stage stays coherent. A future "section_first" jump can be
        # layered on top if intent-shaped queries need it.
        for i, cn in enumerate(chunk_nodes):
            mf = cn.modal_features
            if i > 0:
                mf["preceding_chunk_id"] = {"chunk_id": chunk_ids[i - 1]}
            if i < len(chunk_nodes) - 1:
                mf["following_chunk_id"] = {"chunk_id": chunk_ids[i + 1]}        # Build entity records (UIR v1 doesn't carry per-entity id; the
        # orchestrator keeps the index-based list). Gated on
        # ``include_semantics`` so the agent-facing default JSON stays
        # short -- collecting the entity + relationship lists is cheap
        # (they're already filtered against the boilerplate regex), but
        # PUBLISHING 411 entities + 736 relationships into the default
        # JSON was the real-world complaint that drove the --include-
        # semantics flag. When the flag is off we emit empty lists; the
        # UIR schema validator (Pydantic) accepts the empty default.
        entities: list[Entity] = []
        relationships: list[Relationship] = []
        topics_out: list[str] = []
        if include_semantics:
            entities = [
                Entity(text=e.text, type=e.type, confidence=e.confidence)
                for e in enrichment.entities
            ]
            relationships = [
                Relationship(**{"from": r.from_text},
                              to=r.to_text, type=r.type,
                              confidence=r.confidence)
                for r in enrichment.relationships
            ]
            topics_out = list(enrichment.topics)




        now = datetime.now(timezone.utc)
        provenance = Provenance(
            extraction=ExtractionProvenance(
                model="LayoutLMv3-heuristic",
                version="1.0",
                timestamp=now,
            ),
            normalization=NormalizationProvenance(
                version="1.0",
                timestamp=now,
            ),
        )

        # Lift distinct ``modal_features.section.path`` values into real
        # ``StructureNode(type="section")`` parents (Fix Plan item #3).
        # ``StructureChild`` is a discriminated union so a section node
        # can wrap chunk children with no schema churn. Chunks without a
        # section path stay directly under root. ``current_section``
        # state is local so we always append consecutive same-path chunks
        # to the same parent before opening a new section on path change.
        children: list[Any] = []
        current_section: StructureNode | None = None
        current_path: str = ""
        for cn in chunk_nodes:
            section_path: str = (
                (cn.modal_features.get("section", {}).get("path") or "")
                if cn.modal_features
                else ""
            )
            if section_path:
                if current_section is None or current_path != section_path:
                    sec_id = deterministic_node_id(
                        "section", doc_id, len(children), section_path,
                    )
                    current_section = StructureNode(
                        id=sec_id,
                        type="section",
                        title=section_path,
                        page=cn.page,
                        children=[cn],
                    )
                    children.append(current_section)
                    current_path = section_path
                else:
                    current_section.children.append(cn)
            else:
                current_section = None
                current_path = ""
                children.append(cn)
        root = StructureNode(
            id=doc_id,
            type="document",
            title=metadata.title,
            page=1,
            children=children,
        )
        uir = UIRV1(
            uiR_version="1.0",
            id=doc_id,
            modal_type="document",
            source=source,
            metadata=metadata,
            structure=Structure(type="hierarchical", root=root),
            # Semantics block is gated behind ``include_semantics``. When
            # the flag is off we publish an EMPTY Semantics object so the
            # schema still validates but the JSON payload is hundreds of
            # KB lighter on a real arXiv doc. Topics still get surfaced
            # via UMR's eyebrow / future cross-references; they're cheap.
            semantics=Semantics(
                entities=entities,
                relationships=relationships,
                topics=topics_out,
            ),
            provenance=provenance,
        )

        # Stage 10: write JSON + UMR.
        # UMR (Universal Markdown Representation) is the agent-friendly
        # companion file emitted ALWAYS -- independent of the
        # include_semantics flag. It carries no entities / relationships
        # / topics; only the structured content LLM agents need. See
        # ``src/uir_pipeline/umr.py`` for the rendering contract.
        out_dir = output_dir
        if not dry_run:
            out_dir.mkdir(parents=True, exist_ok=True)
            out_path = out_dir / f"{doc_id}.uir.json"
            # encoding= is mandatory: Path.write_text() defaults to the
            # locale encoding, which is cp1252 on Windows. Any non-latin1
            # glyph in an extracted document (curly quotes, CJK, "·") then
            # either mangles or raises UnicodeEncodeError.
            out_path.write_text(uir.model_dump_json(indent=2), encoding="utf-8")
        else:
            out_path = out_dir / f"{doc_id}.uir.json"  # virtual
        # UMR is rendered from the in-memory UIRV1 (not the just-written
        # JSON file) so a weaviate-networked deployment that disables
        # disk writes for the JSON still gets the agent-facing markdown.
        # ``umr_path`` reflects the path the file WOULD be at on disk;
        # under ``dry_run=True`` the file is not written but the path is
        # still informative for tests / future writers that opt in.
        umr_path = out_dir / f"{doc_id}.umr.md"
        try:
            from uir_pipeline.umr import build_umr
            # UMR consumes the parsed JSON dict shape (not the Pydantic
            # model) so the renderer can be unit-tested independently.
            umr_text = build_umr(json.loads(uir.model_dump_json()))
            if not dry_run:
                umr_path.write_text(umr_text, encoding="utf-8")
        except Exception as exc:
            # UMR rendering is best-effort: if it fails we log and emit a
            # minimal placeholder so downstream consumers can still
            # surface something via the /api/umr/ endpoint instead of
            # 404ing on a missing file.
            logger.warning("UMR render failed (fail-soft): %s", exc)
            if not dry_run:
                umr_path.write_text(
                    f"# UMR render failed\n\n_Exception:_ `{exc}`\n",
                    encoding="utf-8",
                )

        # Stage 11: optional Weaviate upsert
        if not skip_weaviate and not dry_run and vectors is not None and all_chunks:
            client = None
            try:
                from uir_pipeline.weaviate_store import get_client
                client = get_client()
                ensure_collections(client)
                upsert_chunks(client, doc_id, [
                    {
                        "uir_id": cn.id,
                        "text": cn.text,
                        "page": cn.page,
                        "chunk_index": i,
                        "vector": vectors.vectors[i],
                    }
                    for i, cn in enumerate(chunk_nodes)
                ])
                upsert_parent_doc(
                    client, doc_id, mean_pool_vectors(vectors.vectors),
                    extra={"page_count": doc.page_count, "chunk_count": len(chunk_nodes)},
                )
                logger.info("weaviate upsert: %d chunks + 1 doc", len(chunk_nodes))
            except Exception as exc:
                # Fail-soft ONLY when the server isn't there: the CLI defaults
                # to skip_weaviate=False, so a dev without `docker compose up`
                # must still get their UIR JSON. Anything else -- a schema
                # mismatch, a rejected object, a client-API break -- is a bug
                # in us, and swallowing it silently loses the entire index.
                if _is_weaviate_unavailable(exc):
                    logger.warning("weaviate upsert skipped (server unavailable): %s", exc)
                else:
                    logger.error("weaviate upsert failed", exc_info=True)
                    raise
            finally:
                if client is not None:
                    try:
                        client.close()
                    except Exception:  # pragma: no cover -- best-effort cleanup
                        pass

        _progress("done", 100)
        elapsed = time.monotonic() - t0
        return PipelineResult(
            uir_id=doc_id,
            out_path=out_path,
            umr_path=umr_path,
            chunk_count=len(chunk_nodes),
            entity_count=len(entities) if include_semantics else 0,
            elapsed_seconds=round(elapsed, 3),
        )
    finally:
        detach_doc_log(log_handler)


# Re-exports for callers that prefer ``pipeline.derive_doc_id``-style imports.
__all__ = [
    "AudioAnalysisError",
    "ImageAnalysisError",
    "VideoAnalysisError",
    "PipelineResult",
    "run",
]
